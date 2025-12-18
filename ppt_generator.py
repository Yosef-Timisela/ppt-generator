import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import pandas as pd
import io
import uuid

# =============================================
# EXPERT LEVEL PPT GENERATOR (AI-READY & DATA)
# =============================================

st.set_page_config(page_title="Expert PPT Generator", page_icon="🧠", layout="wide")

# ------------------------------
# SESSION STATE
# ------------------------------
if "slides" not in st.session_state:
    st.session_state.slides = []
if "projects" not in st.session_state:
    st.session_state.projects = {}

# ------------------------------
# CUSTOM CSS
# ------------------------------
st.markdown("""
<style>
.main {background-color:#eef2f7;}
.card {background:white;padding:1.2rem;border-radius:16px;box-shadow:0 10px 30px rgba(0,0,0,0.08);margin-bottom:1rem;}
.preview {background-size:cover;background-position:center;border-radius:16px;padding:1.5rem;min-height:200px;box-shadow: inset 0 0 0 2000px rgba(0,0,0,0.35);color:white;}
.stButton>button {background:linear-gradient(135deg,#7c3aed,#4f46e5);color:white;border-radius:12px;font-weight:600;}
</style>
""", unsafe_allow_html=True)

# ------------------------------
# HEADER
# ------------------------------
st.title("🧠 Expert PowerPoint Generator")
st.caption("AI-ready • Data-driven • Project-based • Production architecture")

# ------------------------------
# SIDEBAR – PROJECT MANAGER
# ------------------------------
st.sidebar.header("📂 Project Manager")

project_name = st.sidebar.text_input("Project Name", "My Presentation")

if st.sidebar.button("💾 Save Project"):
    st.session_state.projects[project_name] = st.session_state.slides.copy()
    st.sidebar.success("Project saved")

if st.sidebar.button("📂 Load Project") and project_name in st.session_state.projects:
    st.session_state.slides = st.session_state.projects[project_name].copy()
    st.sidebar.success("Project loaded")

st.sidebar.divider()

# ------------------------------
# SIDEBAR – THEME
# ------------------------------
st.sidebar.header("🎨 Theme Settings")
bg_image = st.sidebar.file_uploader("Background Image", type=["png","jpg","jpeg"])
font_color = st.sidebar.color_picker("Font Color", "#ffffff")

# ------------------------------
# MAIN LAYOUT
# ------------------------------
editor, preview = st.columns([2,1])

# ------------------------------
# SLIDE CREATOR
# ------------------------------
with editor:
    st.subheader("➕ Create Slide")
    slide_type = st.selectbox("Slide Type", ["Text","Bullet","Chart (CSV)"])
    title = st.text_input("Slide Title")
    content = st.text_area("Slide Content / Bullet")
    csv_file = None
    if slide_type == "Chart (CSV)":
        csv_file = st.file_uploader("Upload CSV", type=["csv"])

    if st.button("Add Slide") and title:
        st.session_state.slides.append({
            "id": str(uuid.uuid4()),
            "type": slide_type,
            "title": title,
            "content": content,
            "csv": csv_file
        })

# ------------------------------
# PREVIEW + REORDER
# ------------------------------
with preview:
    st.subheader("👀 Slide Preview")

    for i, slide in enumerate(st.session_state.slides):
        st.markdown(f"""
        <div class="card">
            <div class="preview">
                <h3 style="color:{font_color}">{slide['title']}</h3>
                <p style="color:{font_color}">{slide['content'].replace(chr(10), '<br>')}</p>
                <small>{slide['type']}</small>
            </div>
        </div>
        """, unsafe_allow_html=True)
        col1, col2, col3 = st.columns(3)
        if col1.button("⬆️", key=f"up{slide['id']}") and i>0:
            st.session_state.slides[i-1], st.session_state.slides[i] = st.session_state.slides[i], st.session_state.slides[i-1]
        if col2.button("⬇️", key=f"down{slide['id']}") and i < len(st.session_state.slides)-1:
            st.session_state.slides[i+1], st.session_state.slides[i] = st.session_state.slides[i], st.session_state.slides[i+1]
        if col3.button("🗑️", key=f"del{slide['id']}"):
            st.session_state.slides.pop(i)
            st.rerun()

# ------------------------------
# EXPORT PPT
# ------------------------------
if st.button("🚀 Export PowerPoint"):
    prs = Presentation()

    def add_bg(slide):
        if bg_image:
            slide.shapes.add_picture(io.BytesIO(bg_image.getvalue()), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    for slide_data in st.session_state.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
        tf = box.text_frame
        tf.text = slide_data['title']
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        if slide_data['type'] == "Bullet":
            for line in slide_data['content'].split("\n"):
                p = tf.add_paragraph()
                p.text = line
                p.level = 1
        elif slide_data['type'] == "Chart (CSV)" and slide_data['csv']:
            df = pd.read_csv(slide_data['csv'])
            fig, ax = plt.subplots()
            df.plot(ax=ax)
            img = io.BytesIO()
            plt.savefig(img)
            plt.close()
            img.seek(0)
            slide.shapes.add_picture(img, Inches(1), Inches(2.5), Inches(6))
        else:
            tf.add_paragraph().text = slide_data['content']

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    st.success("🎉 Expert-level PowerPoint created")
    st.download_button("⬇️ Download PPT", buffer, file_name="expert_presentation.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
